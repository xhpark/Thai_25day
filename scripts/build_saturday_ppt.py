#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


ROOT = Path(__file__).resolve().parents[1]
PPT_DIR = ROOT / "assets/generated/ppt"

GREEN = RGBColor(0x2F, 0x7D, 0x5A)
GREEN_DARK = RGBColor(0x17, 0x3F, 0x31)
YELLOW = RGBColor(0xF4, 0xC9, 0x5D)
BLUE = RGBColor(0x17, 0x69, 0xD2)
INK = RGBColor(0x1C, 0x24, 0x1F)
MUTED = RGBColor(0x5D, 0x69, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_BG = RGBColor(0xF4, 0xFB, 0xF4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Malgun Gothic"

TIME_ANCHORS = ["0:00~1:00", "1:00~3:00", "3:00~5:00", "5:00~8:00", "8:00~9:00", "9:00~10:00"]


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def add_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def fill_bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=24, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_rect(slide, left, top, width, height, fill: RGBColor | None, line: RGBColor | None = None, radius=False):
    shape_type = 5 if radius else 1  # ROUNDED_RECTANGLE=5, RECTANGLE=1
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.25)
    shape.shadow.inherit = False
    return shape


def header_band(slide, week: int, day_label: str, title: str, kicker: str):
    fill_bg(slide, RGBColor(0xFB, 0xFA, 0xF5))
    band = add_rect(slide, 0, 0, SLIDE_W, Inches(1.5), GREEN)
    band.fill.fore_color.rgb = GREEN
    tf = band.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"{kicker}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = FONT
    add_text(slide, Inches(9.6), Inches(0.15), Inches(3.4), Inches(1.2),
              day_label, size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, page_no: int, total: int, note: str = ""):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.4), note, size=12, color=MUTED)
    add_text(slide, Inches(11.3), Inches(7.05), Inches(1.6), Inches(0.4), f"{page_no} / {total}",
              size=12, color=MUTED, align=PP_ALIGN.RIGHT)


def slide_opening(prs, spec):
    s = add_slide(prs)
    header_band(s, spec["week"], "토요일 10분 수업", spec["title"], f"{spec['week']}주차 · {spec['theme']}")
    image_path = ROOT / "assets/images/scenes" / f"{spec['primaryImage']['sceneId']}.jpg"
    if image_path.exists():
        s.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.8), height=Inches(3.3))
    add_text(s, Inches(6.6), Inches(1.85), Inches(6.2), Inches(0.5), "오늘의 장면", size=18, bold=True, color=GREEN_DARK)
    add_text(s, Inches(6.6), Inches(2.3), Inches(6.2), Inches(1.6), spec["roleplayScene"], size=18, color=INK, line_spacing=1.3)
    add_rect(s, Inches(6.6), Inches(4.1), Inches(6.2), Inches(1.6), PANEL_BG, RGBColor(0xD7, 0xEA, 0xD9), radius=True)
    add_text(s, Inches(6.85), Inches(4.25), Inches(5.7), Inches(0.4), "수업 목표", size=16, bold=True, color=GREEN_DARK)
    add_text(s, Inches(6.85), Inches(4.65), Inches(5.7), Inches(0.9), spec["missionGoal"], size=17, color=INK, line_spacing=1.25)
    footer(s, 1, 8, "1주차 토요일 오프라인 수업 · 태국어 선교회화 25일")
    return s


def slide_core_phrases(prs, spec):
    s = add_slide(prs)
    header_band(s, spec["week"], "10분 수업", "이번 주 핵심 문장 전체", "Slide 2 · 핵심 문장")
    phrases = spec["phrases"]
    cols = 3
    gap = Inches(0.18)
    col_w = Inches(4.0)
    rows = -(-len(phrases) // cols)  # ceil
    row_h = Inches((6.7 - 1.75 - (rows - 1) * 0.18) / rows)
    start_top = Inches(1.75)
    for idx, phrase in enumerate(phrases):
        col = idx % cols
        row = idx // cols
        left = Inches(0.55) + col * (col_w + gap)
        top = start_top + row * (row_h + gap)
        card = add_rect(s, left, top, col_w, row_h, WHITE, RGBColor(0xDF, 0xE5, 0xDC), radius=True)
        tf = card.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phrase["korean"]
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = GREEN_DARK
        p.font.name = FONT
        female = phrase["speech"]["female"]
        male = phrase["speech"]["male"]
        p2 = tf.add_paragraph()
        p2.text = female["thai"]
        p2.font.size = Pt(23)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0x05, 0x05, 0x05)
        p2.font.name = FONT
        p3 = tf.add_paragraph()
        p3.text = f"여: {female['korean_pronunciation']}"
        p3.font.size = Pt(13)
        p3.font.bold = True
        p3.font.color.rgb = BLUE
        p3.font.name = FONT
        p4 = tf.add_paragraph()
        p4.text = f"남: {male['korean_pronunciation']}"
        p4.font.size = Pt(13)
        p4.font.bold = True
        p4.font.color.rgb = BLUE
        p4.font.name = FONT
    footer(s, 2, 8, "여성 발화 먼저 연습 → 남성형 비교")
    return s


def slide_keywords(prs, spec):
    s = add_slide(prs)
    header_band(s, spec["week"], "10분 수업", "핵심 단어", "Slide 3 · 핵심 단어 + 이미지")
    keywords = spec["keywords"][:4]
    n = max(len(keywords), 1)
    card_w = Inches(11.8 / n)
    top = Inches(2.2)
    height = Inches(3.4)
    for idx, kw in enumerate(keywords):
        left = Inches(0.7) + idx * card_w
        card = add_rect(s, left, top, card_w - Inches(0.2), height, PANEL_BG, RGBColor(0xD7, 0xEA, 0xD9), radius=True)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        p = tf.paragraphs[0]
        p.text = kw["thai"]
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = INK
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = kw["koreanPronunciation"]
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = BLUE
        p2.font.name = FONT
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = kw["korean"]
        p3.font.size = Pt(18)
        p3.font.color.rgb = GREEN_DARK
        p3.font.name = FONT
        p3.alignment = PP_ALIGN.CENTER
        image_path = ROOT / "assets/images/keywords" / f"{kw['id']}.jpg"
        if image_path.exists():
            s.shapes.add_picture(str(image_path), left + Inches(0.3), top + Inches(1.6), width=card_w - Inches(0.8))
        else:
            add_text(s, left, top + Inches(2.0), card_w - Inches(0.2), Inches(1.2),
                      "(단어 이미지 준비 중)", size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 3, 8, "이미지가 없는 단어는 음성과 발화로만 연습합니다.")
    return s


def slide_male_female(prs, spec):
    s = add_slide(prs)
    header_band(s, spec["week"], "10분 수업", "남성형 / 여성형 어미 점검", "Slide 4 · 발화 비교")
    headers = ["문장", "남성형 (ครับ)", "여성형 (ค่ะ)"]
    col_w = [Inches(3.6), Inches(4.3), Inches(4.3)]
    left0 = Inches(0.6)
    top = Inches(1.7)
    row_h = Inches(0.6)
    x = left0
    for i, h in enumerate(headers):
        hdr = add_rect(s, x, top, col_w[i], row_h, GREEN)
        tf = hdr.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = h
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        x += col_w[i]
    phrases = spec["phrases"]
    for row_idx, phrase in enumerate(phrases):
        ytop = top + row_h * (row_idx + 1)
        x = left0
        male = phrase["speech"]["male"]
        female = phrase["speech"]["female"]
        cells = [phrase["korean"], f"{male['thai']}\n{male['korean_pronunciation']}", f"{female['thai']}\n{female['korean_pronunciation']}"]
        for i, text in enumerate(cells):
            fill = WHITE if row_idx % 2 == 0 else RGBColor(0xF4, 0xFB, 0xF4)
            cell = add_rect(s, x, ytop, col_w[i], row_h, fill, RGBColor(0xDF, 0xE5, 0xDC))
            tf = cell.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            p = tf.paragraphs[0]
            lines = text.split("\n")
            p.text = lines[0]
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.name = FONT
            p.font.color.rgb = INK
            p.alignment = PP_ALIGN.CENTER
            if len(lines) > 1:
                p2 = tf.add_paragraph()
                p2.text = lines[1]
                p2.font.size = Pt(13)
                p2.font.color.rgb = BLUE
                p2.font.name = FONT
                p2.alignment = PP_ALIGN.CENTER
            x += col_w[i]
    table_bottom = top + row_h * (len(phrases) + 1)
    add_text(s, Inches(0.6), table_bottom + Inches(0.06), Inches(11.2), Inches(0.4),
              "포인트: 남성 학습자는 끝에 '크랍', 여성 학습자는 '카/케'로 마무리합니다.", size=14, bold=True, color=GREEN_DARK)
    footer(s, 4, 8)
    return s


def slide_roleplay_steps(prs, spec, teacher_flow: list[str]):
    s = add_slide(prs)
    header_band(s, spec["week"], "10분 수업", "토요일 10분 진행 순서", "Slide 5 · 역할극 단계")
    top = Inches(1.85)
    row_h = Inches(0.78)
    for idx, step in enumerate(teacher_flow):
        y = top + idx * (row_h + Inches(0.06))
        time_label = TIME_ANCHORS[idx] if idx < len(TIME_ANCHORS) else ""
        pill = add_rect(s, Inches(0.6), y, Inches(1.6), row_h, GREEN, radius=True)
        tf = pill.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = time_label
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        card = add_rect(s, Inches(2.35), y, Inches(10.4), row_h, WHITE, RGBColor(0xDF, 0xE5, 0xDC), radius=True)
        tf2 = card.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf2.margin_left = Inches(0.25)
        p2 = tf2.paragraphs[0]
        p2.text = f"{idx + 1}. {step}"
        p2.font.size = Pt(19)
        p2.font.bold = True
        p2.font.color.rgb = INK
        p2.font.name = FONT
    footer(s, 5, 8, "역할극 핵심 표현: " + ", ".join(p["korean"] for p in spec["phrases"]))
    return s


def slide_ministry_guide(prs, spec):
    s = add_slide(prs)
    header_band(s, spec["week"], "10분 수업", "선교 현장 가이드", "Slide 6 · 미션 가이드")
    add_rect(s, Inches(0.6), Inches(1.9), Inches(11.8), Inches(2.0), RGBColor(0xED, 0xF7, 0xFF), RGBColor(0xC9, 0xDF, 0xF1), radius=True)
    add_text(s, Inches(0.9), Inches(2.05), Inches(11.2), Inches(0.4), "강사 가이드", size=18, bold=True, color=GREEN_DARK)
    add_text(s, Inches(0.9), Inches(2.5), Inches(11.2), Inches(1.3), spec["ministryGuide"], size=19, color=INK, line_spacing=1.35)
    add_rect(s, Inches(0.6), Inches(4.1), Inches(11.8), Inches(1.9), RGBColor(0xFF, 0xF8, 0xDF), RGBColor(0xEF, 0xD7, 0x7B), radius=True)
    add_text(s, Inches(0.9), Inches(4.25), Inches(11.2), Inches(0.4), "성공의 모습", size=18, bold=True, color=RGBColor(0x8A, 0x66, 0x14))
    add_text(s, Inches(0.9), Inches(4.7), Inches(11.2), Inches(1.2),
              "학습자가 남성형/여성형 어미를 구분하며 핵심 표현을 짝과 함께 소리 내어 말하고, 따뜻한 선교적 태도로 역할극에 참여하면 성공이다.",
              size=18, color=INK, line_spacing=1.3)
    footer(s, 6, 8)
    return s


def slide_practice_round(prs, spec):
    s = add_slide(prs)
    header_band(s, spec["week"], "10분 수업", "짝 역할극 연습", "Slide 7 · 연습 라운드")
    add_rect(s, Inches(0.6), Inches(1.9), Inches(5.7), Inches(4.1), WHITE, RGBColor(0xDF, 0xE5, 0xDC), radius=True)
    add_text(s, Inches(0.9), Inches(2.1), Inches(5.1), Inches(0.5), "1라운드", size=20, bold=True, color=GREEN_DARK)
    add_text(s, Inches(0.9), Inches(2.7), Inches(5.1), Inches(3.0),
              "두 명씩 짝을 지어 오늘 장면을 처음부터 끝까지 한 번 말해 봅니다.\n\n"
              "한 사람이 방문자, 한 사람이 현지 봉사자 역할을 맡습니다.\n\n"
              "발음이 완벽하지 않아도 끝까지 말하는 것이 목표입니다.",
              size=17, color=INK, line_spacing=1.35)
    add_rect(s, Inches(6.6), Inches(1.9), Inches(5.7), Inches(4.1), WHITE, RGBColor(0xDF, 0xE5, 0xDC), radius=True)
    add_text(s, Inches(6.9), Inches(2.1), Inches(5.1), Inches(0.5), "2라운드 (역할 교체)", size=20, bold=True, color=GREEN_DARK)
    add_text(s, Inches(6.9), Inches(2.7), Inches(5.1), Inches(3.0),
              "역할을 바꾸어 같은 장면을 한 번 더 말해 봅니다.\n\n"
              "이번에는 조금 더 자신 있게, 눈을 마주치며 말해 봅니다.\n\n"
              "짝과 함께 잘한 점을 한 가지씩 칭찬해 줍니다.",
              size=17, color=INK, line_spacing=1.35)
    footer(s, 7, 8)
    return s


def slide_sendoff(prs, spec, sunday_mission: str | None):
    s = add_slide(prs)
    header_band(s, spec["week"], "10분 수업", "마무리와 다음 주 예고", "Slide 8 · 보내기")
    add_rect(s, Inches(0.6), Inches(1.9), Inches(11.8), Inches(1.8), PANEL_BG, RGBColor(0xD7, 0xEA, 0xD9), radius=True)
    add_text(s, Inches(0.9), Inches(2.05), Inches(11.2), Inches(0.4), "이번 주 잘한 점 칭찬하기", size=18, bold=True, color=GREEN_DARK)
    add_text(s, Inches(0.9), Inches(2.5), Inches(11.2), Inches(1.1),
              "발음보다 먼저 입을 연 것, 눈을 맞추고 미소로 인사한 것을 함께 칭찬합니다.",
              size=18, color=INK, line_spacing=1.3)
    add_rect(s, Inches(0.6), Inches(3.9), Inches(11.8), Inches(2.1), RGBColor(0xFF, 0xF8, 0xDF), RGBColor(0xEF, 0xD7, 0x7B), radius=True)
    add_text(s, Inches(0.9), Inches(4.05), Inches(11.2), Inches(0.4), "일요일 개인 활용 미션", size=18, bold=True, color=RGBColor(0x8A, 0x66, 0x14))
    add_text(s, Inches(0.9), Inches(4.5), Inches(11.2), Inches(1.4),
              sunday_mission or "이번 주 배운 표현을 가족이나 가까운 사람에게 직접 말해 봅니다.",
              size=18, color=INK, line_spacing=1.3)
    footer(s, 8, 8, "다음 주에 또 만나요! 태국어 선교회화 25일")
    return s


def build(lesson_id: str) -> Path:
    ppt_spec = load_json(PPT_DIR / f"{lesson_id.lower()}.json")
    script_spec = load_json(PPT_DIR / f"{lesson_id.lower()}_script.json")
    teacher_flow = script_spec.get("teacherFlow", [])

    sunday_id = f"w{ppt_spec['week']}sun"
    sunday_mission = None
    sunday_path = ROOT / "assets/generated/print" / f"{sunday_id}.json"
    if sunday_path.exists():
        sunday_spec = load_json(sunday_path)
        sunday_mission = sunday_spec.get("personalMission")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_opening(prs, ppt_spec)
    slide_core_phrases(prs, ppt_spec)
    slide_keywords(prs, ppt_spec)
    slide_male_female(prs, ppt_spec)
    slide_roleplay_steps(prs, ppt_spec, teacher_flow)
    slide_ministry_guide(prs, ppt_spec)
    slide_practice_round(prs, ppt_spec)
    slide_sendoff(prs, ppt_spec, sunday_mission)

    out_path = PPT_DIR / f"{lesson_id.lower()}.pptx"
    prs.save(str(out_path))
    return out_path


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a Saturday offline-lesson PPTX from a saturday_ppt_spec JSON.")
    parser.add_argument("lesson_id", nargs="?", default="w1sat", help="e.g. w1sat, w2sat")
    args = parser.parse_args()
    out_path = build(args.lesson_id)
    print(f"[OK] wrote {out_path.relative_to(ROOT)}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
